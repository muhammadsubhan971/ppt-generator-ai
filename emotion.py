from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.oxml.xmlchemy import OxmlElement
from pptx.oxml.ns import nsdecls, qn
import re
import logging
import html
import os
from typing import List, Dict, Tuple, Optional, Union


class EnhancedPPTGenerator:
    """Enhanced PowerPoint presentation generator with template support instead of theme colors."""
    
    # Class constants - IMPROVED for better content handling
    MAX_BULLETS_PER_SLIDE = 5  # Reduced for more comprehensive bullets
    DEFAULT_TITLE_FONT_SIZE = 44
    DEFAULT_SUBTITLE_FONT_SIZE = 24
    DEFAULT_CONTENT_FONT_SIZE = 16  # Slightly smaller for comprehensive content
    SECTION_TITLE_FONT_SIZE = 36
    COMPREHENSIVE_TEXT_THRESHOLD = 120  # Increased threshold for detailed content
    MAX_TITLE_LENGTH = 65
    
    def __init__(self, template_name: str = "default"):
        """Initialize the PPT generator with specified template."""
        try:
            self.template_name = template_name
            self.template_path = self._get_template_path(template_name)
            
            # Load template or create new presentation
            if self.template_path and os.path.exists(self.template_path):
                self.ppt = Presentation(self.template_path)
                print(f"✅ Loaded template: {template_name}")
            else:
                self.ppt = Presentation()
                print(f"⚠️ Template '{template_name}' not found. Using default PowerPoint template.")
            
            self._setup_layouts()
            self._setup_logging()
            logging.info(f"Initialized PPT Generator with template: {template_name}")
        except Exception as e:
            logging.error(f"Failed to initialize PPTGenerator: {e}")
            raise
    
    def _get_template_path(self, template_name: str) -> Optional[str]:
        """Get the full path to the template file."""
        if template_name == "default" or not template_name:
            return None
        
        templates_folder = "templates"
        
        # If template_name already has .pptx extension
        if template_name.lower().endswith('.pptx'):
            template_path = os.path.join(templates_folder, template_name)
        else:
            template_path = os.path.join(templates_folder, f"{template_name}.pptx")
        
        return template_path if os.path.exists(template_path) else None
    
    def _setup_layouts(self):
        """Setup slide layouts with error handling."""
        try:
            # Get available layouts from the presentation
            available_layouts = len(self.ppt.slide_layouts)
            
            self.title_slide_layout = self.ppt.slide_layouts[0]
            self.title_content_layout = self.ppt.slide_layouts[1] if available_layouts > 1 else self.ppt.slide_layouts[0]
            self.section_layout = (self.ppt.slide_layouts[2] 
                                 if available_layouts > 2 
                                 else self.ppt.slide_layouts[1] if available_layouts > 1 
                                 else self.ppt.slide_layouts[0])
            
            logging.info(f"Slide layouts configured successfully. Available layouts: {available_layouts}")
        except IndexError as e:
            logging.error(f"Layout setup failed: {e}")
            raise
    
    def _setup_logging(self):
        """Setup logging configuration."""
        logging.basicConfig(
            level=logging.INFO,
            format='%(asctime)s - %(levelname)s - %(message)s'
        )
    
    def _clean_html_entities(self, text: str) -> str:
        """Clean HTML entities and malformed encoding."""
        if not text:
            return ""
        
        # Decode HTML entities like &amp;, &lt;, &gt;, etc.
        text = html.unescape(text)
        
        # Fix common HTML entity issues
        replacements = {
            '&amp;': '&',
            '&lt;': '<',
            '&gt;': '>',
            '&nbsp;': ' ',
            '&quot;': '"',
            '&#39;': "'",
            '&hellip;': '...',
            '&mdash;': '—',
            '&ndash;': '–'
        }
        
        for old, new in replacements.items():
            text = text.replace(old, new)
        
        return text
    
    def _clean_escape_characters(self, text: str) -> str:
        """Remove unnecessary escape characters and fix formatting."""
        if not text:
            return ""
        
        # Remove unnecessary backslashes before punctuation
        text = re.sub(r'\\([.,;:!?()[\]{}])', r'\1', text)
        
        # Fix escaped formatting characters
        text = re.sub(r'\\\*', '*', text)
        text = re.sub(r'\\_', '_', text)
        text = re.sub(r'\\-', '-', text)
        text = re.sub(r'\\"', '"', text)
        text = re.sub(r"\\'", "'", text)
        
        return text
    
    def _remove_html_tags(self, text: str) -> str:
        """Remove HTML tags while preserving content."""
        if not text:
            return ""
        
        # Remove span tags but keep content
        text = re.sub(r'<span[^>]*>(.*?)</span>', r'\1', text, flags=re.DOTALL)
        
        # Remove other common HTML tags
        text = re.sub(r'</?(?:div|p|br|strong|em|b|i|u|font)[^>]*>', '', text)
        
        # Clean up any remaining HTML tags
        text = re.sub(r'<[^>]+>', '', text)
        
        return text
    
    def _normalize_whitespace(self, text: str) -> str:
        """Normalize whitespace and line breaks."""
        if not text:
            return ""
        
        # Replace multiple spaces with single space
        text = re.sub(r'\s+', ' ', text)
        
        # Remove leading/trailing whitespace
        text = text.strip()
        
        return text
    
    def _validate_and_clean_text_input(self, text: Union[str, None]) -> str:
        """Comprehensive text validation and cleaning."""
        if text is None:
            return ""
        if not isinstance(text, str):
            text = str(text)
        
        # Apply all cleaning functions in sequence
        text = self._clean_html_entities(text)
        text = self._clean_escape_characters(text)
        text = self._remove_html_tags(text)
        text = self._normalize_whitespace(text)
        
        return text
    
    def _process_text_formatting(self, text: str) -> str:
        """Process markdown-style formatting in text with improved handling."""
        if not text:
            return ""
        
        # Clean the text first
        text = self._validate_and_clean_text_input(text)
        
        # Handle bold formatting: **text** -> preserve for later processing
        text = re.sub(r'\*\*\s*([^*]+?)\s*\*\*', r'**\1**', text)
        
        # Handle italic formatting: *text* -> preserve for later processing
        text = re.sub(r'(?<!\*)\*\s*([^*]+?)\s*\*(?!\*)', r'*\1*', text)
        
        # Clean up other markdown-style formatting
        text = re.sub(r'__\s*([^_]+?)\s*__', r'**\1**', text)  # Convert __ to **
        text = re.sub(r'~~\s*([^~]+?)\s*~~', r'\1', text)      # Remove strikethrough
        
        return text
    
    def _apply_text_formatting(self, paragraph, original_text: str):
        """Apply rich text formatting with improved parsing."""
        if not original_text:
            return paragraph
        
        # Clean the text first
        text = self._process_text_formatting(original_text)
        
        # Clear existing text
        paragraph.clear()
        
        # Split text by bold markers and process
        parts = re.split(r'(\*\*[^*]+?\*\*)', text)
        
        if not parts or all(not part.strip() for part in parts):
            # Fallback: add original text as single run
            run = paragraph.add_run()
            run.text = original_text
            run.font.size = Pt(self.DEFAULT_CONTENT_FONT_SIZE)
            run.font.name = "Calibri"
            return paragraph
        
        for part in parts:
            if not part:
                continue
            
            run = paragraph.add_run()
            
            # Check if this part is bold
            if part.startswith('**') and part.endswith('**') and len(part) > 4:
                run.text = part[2:-2]  # Remove ** markers
                run.font.bold = True
            else:
                # Handle italic text within non-bold parts
                italic_parts = re.split(r'(\*[^*]+?\*)', part)
                
                for i, italic_part in enumerate(italic_parts):
                    if not italic_part:
                        continue
                    
                    if i > 0:  # Not the first part, need new run
                        run = paragraph.add_run()
                    
                    if (italic_part.startswith('*') and italic_part.endswith('*') 
                        and len(italic_part) > 2 and not italic_part.startswith('**')):
                        run.text = italic_part[1:-1]  # Remove * markers
                        run.font.italic = True
                    else:
                        run.text = italic_part
            
            # Apply consistent formatting
            run.font.size = Pt(self.DEFAULT_CONTENT_FONT_SIZE)
            run.font.name = "Calibri"
        
        return paragraph
    
    def _estimate_text_length(self, text: str) -> bool:
        """IMPROVED: Better text length estimation for comprehensive content."""
        if not text:
            return False
        
        # Clean text first for accurate measurement
        clean_text = self._validate_and_clean_text_input(text)
        
        # Consider both character count and word count - ADJUSTED for comprehensive content
        char_count = len(clean_text)
        word_count = len(clean_text.split())
        
        # Allow longer content since we want comprehensive bullets
        return char_count > self.COMPREHENSIVE_TEXT_THRESHOLD or word_count > 20
    
    def _split_long_bullet(self, text: str) -> List[str]:
        """IMPROVED: Better handling of comprehensive content - only split if absolutely necessary."""
        if not text:
            return [text]
        
        # Clean text first
        text = self._validate_and_clean_text_input(text)
        
        # For comprehensive content, be more conservative about splitting
        # Only split if content is extremely long (more than 250 characters)
        if len(text) <= 250:
            return [text]  # Keep comprehensive content intact
        
        # Try splitting at sentence boundaries first
        sentences = re.split(r'(?<=[.!?])\s+', text)
        if len(sentences) > 1:
            # Group sentences to maintain comprehensive content
            grouped_sentences = []
            current_group = []
            current_length = 0
            
            for sentence in sentences:
                sentence = sentence.strip()
                if current_length + len(sentence) <= 200 and current_group:
                    current_group.append(sentence)
                    current_length += len(sentence)
                else:
                    if current_group:
                        grouped_sentences.append(' '.join(current_group))
                    current_group = [sentence]
                    current_length = len(sentence)
            
            if current_group:
                grouped_sentences.append(' '.join(current_group))
            
            return [s for s in grouped_sentences if s.strip()]
        
        # Only as last resort, split at clause boundaries
        if len(text) > 300:  # Very long content
            clauses = re.split(r'(?<=[,;])\s+', text)
            if len(clauses) > 1:
                # Group clauses intelligently
                grouped_clauses = []
                current_group = []
                current_length = 0
                
                for clause in clauses:
                    clause = clause.strip()
                    if current_length + len(clause) <= 180 and current_group:
                        current_group.append(clause)
                        current_length += len(clause)
                    else:
                        if current_group:
                            grouped_clauses.append(', '.join(current_group))
                        current_group = [clause]
                        current_length = len(clause)
                
                if current_group:
                    grouped_clauses.append(', '.join(current_group))
                
                return [s for s in grouped_clauses if s.strip()]
        
        return [text]  # Keep as single comprehensive bullet
    
    def _truncate_title(self, title: str) -> str:
        """Truncate title if too long while preserving meaning."""
        title = self._validate_and_clean_text_input(title)
        
        if len(title) <= self.MAX_TITLE_LENGTH:
            return title
        
        # Try to truncate at word boundary
        words = title.split()
        truncated = ""
        
        for word in words:
            if len(truncated + " " + word) <= self.MAX_TITLE_LENGTH - 3:
                truncated += (" " + word) if truncated else word
            else:
                break
        
        return truncated + "..." if truncated != title else title[:self.MAX_TITLE_LENGTH-3] + "..."
    
    def add_title_slide(self, title: str, subtitle: Optional[str] = None):
        """Add enhanced title slide using template layout."""
        try:
            slide = self.ppt.slides.add_slide(self.title_slide_layout)
            
            # Clean and validate inputs
            title = self._validate_and_clean_text_input(title)
            subtitle = self._validate_and_clean_text_input(subtitle) if subtitle else None
            
            # Set title with proper formatting
            if slide.shapes.title:
                title_shape = slide.shapes.title
                title_shape.text = title
                
                # Configure text frame for better handling
                text_frame = title_shape.text_frame
                text_frame.word_wrap = True
                text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
                text_frame.margin_left = Inches(0.2)
                text_frame.margin_right = Inches(0.2)
                text_frame.margin_top = Inches(0.1)
                text_frame.margin_bottom = Inches(0.1)
                
                # Apply formatting (let template handle colors)
                if text_frame.paragraphs:
                    para = text_frame.paragraphs[0]
                    para.alignment = PP_ALIGN.CENTER
                    
                    # Apply to all runs in paragraph
                    for run in para.runs:
                        run.font.size = Pt(self.DEFAULT_TITLE_FONT_SIZE)
                        run.font.bold = True
                        run.font.name = "Calibri"
            
            # Set subtitle
            if subtitle and len(slide.placeholders) > 1:
                subtitle_shape = slide.placeholders[1]
                subtitle_shape.text = subtitle
                
                text_frame = subtitle_shape.text_frame
                text_frame.word_wrap = True
                text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
                
                if text_frame.paragraphs:
                    para = text_frame.paragraphs[0]
                    para.alignment = PP_ALIGN.CENTER
                    
                    # Apply to all runs in paragraph
                    for run in para.runs:
                        run.font.size = Pt(self.DEFAULT_SUBTITLE_FONT_SIZE)
                        run.font.name = "Calibri"
            
            logging.info(f"Title slide created successfully with template: {self.template_name}")
            return slide
            
        except Exception as e:
            logging.error(f"Failed to create title slide: {e}")
            raise
    
    def add_section_slide(self, title: str, content: List[str], 
                         slide_number: int = 1, total_slides: int = 1):
        """IMPROVED: Add enhanced section slide with better handling of comprehensive content."""
        try:
            slide = self.ppt.slides.add_slide(self.title_content_layout)
            
            # Clean and validate title
            title = self._validate_and_clean_text_input(title)
            title = self._truncate_title(title)
            
            # Clean and validate content - IMPROVED for comprehensive content
            if not isinstance(content, list):
                content = [str(content)] if content else []
            
            cleaned_content = []
            for item in content:
                cleaned_item = self._validate_and_clean_text_input(str(item))
                # CHANGED: Accept shorter content too since AI should generate comprehensive content
                if cleaned_item and len(cleaned_item) > 10:  # Reduced minimum length
                    cleaned_content.append(cleaned_item)
            
            # Set slide title
            if slide.shapes.title:
                title_shape = slide.shapes.title
                title_shape.text = title
                
                text_frame = title_shape.text_frame
                text_frame.word_wrap = True
                text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
                text_frame.margin_left = Inches(0.3)
                text_frame.margin_right = Inches(0.3)
                
                if text_frame.paragraphs:
                    para = text_frame.paragraphs[0]
                    # Ensure we have at least one run
                    if not para.runs:
                        run = para.add_run()
                        run.text = title
                    
                    for run in para.runs:
                        run.font.size = Pt(self.SECTION_TITLE_FONT_SIZE)
                        run.font.bold = True
                        run.font.name = "Calibri"
            
            # IMPROVED: Add content with better formatting for comprehensive bullets
            if len(slide.placeholders) > 1 and cleaned_content:
                content_shape = slide.placeholders[1]
                text_frame = content_shape.text_frame
                text_frame.clear()
                text_frame.word_wrap = True
                text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
                text_frame.margin_left = Inches(0.3)
                text_frame.margin_right = Inches(0.3)
                text_frame.margin_top = Inches(0.1)
                text_frame.margin_bottom = Inches(0.1)
                
                # Add each bullet point with proper formatting
                for i, point in enumerate(cleaned_content):
                    if i >= self.MAX_BULLETS_PER_SLIDE:
                        break
                    
                    # Create paragraph for each bullet point
                    if i == 0:
                        p = text_frame.paragraphs[0]
                    else:
                        p = text_frame.add_paragraph()
                    
                    p.level = 0
                    p.space_after = Pt(12)  # INCREASED space for comprehensive content
                    
                    # IMPROVED: Better text formatting for comprehensive content
                    self._apply_text_formatting(p, point)
                    
                    # Ensure proper formatting is applied
                    for run in p.runs:
                        if not run.font.size:
                            run.font.size = Pt(self.DEFAULT_CONTENT_FONT_SIZE)
                        if not run.font.name:
                            run.font.name = "Calibri"
            
            logging.info(f"Section slide created: {title} ({slide_number}/{total_slides})")
            return slide
            
        except Exception as e:
            logging.error(f"Failed to create section slide: {e}")
            raise
    
    def _distribute_content(self, title: str, content: List[str], 
                          max_slides: Optional[int] = None) -> List[Tuple[str, List[str]]]:
        """IMPROVED: Enhanced content distribution for comprehensive bullets."""
        if not content:
            return [(title, [])]
        
        # Process and clean all content - IMPROVED for comprehensive content
        processed_content = []
        for point in content:
            cleaned_point = self._validate_and_clean_text_input(str(point))
            if not cleaned_point:
                continue
            
            # CHANGED: Be more conservative about splitting comprehensive content
            if self._estimate_text_length(cleaned_point):
                # Only split if absolutely necessary
                split_points = self._split_long_bullet(cleaned_point)
                processed_content.extend(split_points)
            else:
                processed_content.append(cleaned_point)
        
        if not processed_content:
            return [(title, [])]
        
        # IMPROVED: Better distribution for comprehensive content
        if max_slides and max_slides > 0:
            # Fewer bullets per slide for comprehensive content
            points_per_slide = max(1, min(self.MAX_BULLETS_PER_SLIDE, 
                                        (len(processed_content) + max_slides - 1) // max_slides))
        else:
            points_per_slide = self.MAX_BULLETS_PER_SLIDE
        
        # Distribute content
        slides_content = []
        for i in range(0, len(processed_content), points_per_slide):
            slide_content = processed_content[i:i + points_per_slide]
            slide_title = title
            
            # Add part indicator for multi-slide sections
            if len(processed_content) > points_per_slide:
                part_num = (i // points_per_slide) + 1
                total_parts = (len(processed_content) + points_per_slide - 1) // points_per_slide
                slide_title = f"{title} (Part {part_num}/{total_parts})"
            
            slides_content.append((slide_title, slide_content))
        
        return slides_content
    
    def generate_from_content(self, content: Dict) -> Tuple[Presentation, int]:
        """Generate enhanced PowerPoint with template support and comprehensive content handling."""
        try:
            if not isinstance(content, dict):
                raise ValueError("Content must be a dictionary")
            
            # Clean and validate all input content
            title = self._validate_and_clean_text_input(content.get("title", "Presentation"))
            subtitle = self._validate_and_clean_text_input(content.get("subtitle", ""))
            target_slides = max(1, int(content.get("target_slides", 15)))
            
            logging.info(f"Generating presentation with template: {self.template_name}")
            
            # Add title slide
            self.add_title_slide(title, subtitle if subtitle else None)
            
            # Process sections
            sections = content.get("sections", [])
            if not sections:
                logging.warning("No sections found in content")
                return self.ppt, len(self.ppt.slides)
            
            # IMPROVED: Better slide distribution for comprehensive content
            total_content_items = sum(len(section.get("content", [])) for section in sections)
            available_slides = max(1, target_slides - 2)  # Reserve for title and closing
            
            # Generate slides
            slide_count = 1  # Start from 1 (title slide)
            for section_idx, section in enumerate(sections):
                section_title = self._validate_and_clean_text_input(section.get("title", f"Section {section_idx + 1}"))
                section_content = section.get("content", [])
                
                if not section_content:
                    logging.warning(f"No content found for section: {section_title}")
                    continue
                
                # IMPROVED: Better distribution for comprehensive content
                max_slides_for_section = max(1, available_slides // len(sections))
                distributed_content = self._distribute_content(
                    section_title, section_content, max_slides_for_section
                )
                
                # Create slides
                for slide_title, slide_content in distributed_content:
                    slide_count += 1
                    self.add_section_slide(slide_title, slide_content, slide_count, target_slides)
            
            # Add closing slide
            call_to_action = self._validate_and_clean_text_input(content.get("call_to_action", ""))
            self.add_closing_slide("Thank You", call_to_action if call_to_action else None)
            
            actual_slides = len(self.ppt.slides)
            logging.info(f"Presentation generated successfully: {actual_slides} slides with template {self.template_name}")
            
            return self.ppt, actual_slides
            
        except Exception as e:
            logging.error(f"Failed to generate presentation: {e}")
            raise
    
    def add_closing_slide(self, title: str = "Thank You", content: Optional[str] = None):
        """Add enhanced closing slide using template."""
        try:
            slide = self.ppt.slides.add_slide(self.title_slide_layout)
            
            title = self._validate_and_clean_text_input(title)
            content = self._validate_and_clean_text_input(content) if content else None
            
            # Set title
            if slide.shapes.title:
                title_shape = slide.shapes.title
                title_shape.text = title
                
                text_frame = title_shape.text_frame
                text_frame.word_wrap = True
                text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
                
                if text_frame.paragraphs:
                    para = text_frame.paragraphs[0]
                    para.alignment = PP_ALIGN.CENTER
                    
                    # Ensure we have at least one run
                    if not para.runs:
                        run = para.add_run()
                        run.text = title
                    
                    for run in para.runs:
                        run.font.size = Pt(52)
                        run.font.bold = True
                        run.font.name = "Calibri"
            
            # Add content if provided
            if content and len(slide.placeholders) > 1:
                content_shape = slide.placeholders[1]
                content_shape.text = content
                
                text_frame = content_shape.text_frame
                text_frame.word_wrap = True
                text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
                
                if text_frame.paragraphs:
                    para = text_frame.paragraphs[0]
                    para.alignment = PP_ALIGN.CENTER
                    
                    # Ensure we have at least one run
                    if not para.runs:
                        run = para.add_run()
                        run.text = content
                    
                    for run in para.runs:
                        run.font.size = Pt(28)
                        run.font.name = "Calibri"
            
            logging.info(f"Closing slide created with template {self.template_name}")
            return slide
            
        except Exception as e:
            logging.error(f"Failed to create closing slide: {e}")
            raise
    
    def save(self, filename: str = "enhanced_presentation.pptx") -> str:
        """Save presentation with validation and proper naming."""
        try:
            if not filename:
                filename = f"presentation_{self.template_name}.pptx"
            
            if not filename.lower().endswith('.pptx'):
                filename += '.pptx'
            
            # Add template name to filename if not already present
            if self.template_name != "default" and self.template_name not in filename:
                name_part = filename.replace('.pptx', '')
                filename = f"{name_part}_{self.template_name}.pptx"
            
            self.ppt.save(filename)
            logging.info(f"Presentation saved as: {filename}")
            
            return filename
            
        except Exception as e:
            logging.error(f"Failed to save presentation: {e}")
            raise

    def get_template_info(self) -> Dict:
        """Get information about the current template."""
        return {
            "template_name": self.template_name,
            "template_path": self.template_path,
            "total_layouts": len(self.ppt.slide_layouts),
            "template_exists": self.template_path and os.path.exists(self.template_path) if self.template_path else False
        }


# IMPROVED: Enhanced example usage and testing
if __name__ == "__main__":
    # Create templates folder if it doesn't exist
    templates_folder = "templates"
    if not os.path.exists(templates_folder):
        os.makedirs(templates_folder)
        print(f"Created {templates_folder} folder. Please add your .pptx template files there.")
    
    # Comprehensive sample content with detailed bullet points
    sample_content = {
        "title": "Advanced Academic Presentation",
        "subtitle": "Enhanced with Template Support",
        "target_slides": 15,
        "sections": [
            {
                "title": "Introduction to Complex Systems",
                "content": [
                    "Complex systems are characterized by multiple interacting components that exhibit emergent behaviors not predictable from individual component properties alone",
                    "These systems demonstrate non-linear dynamics where small changes can lead to disproportionately large effects throughout the entire system structure",
                    "Emergent properties arise from the collective behavior of system components and cannot be understood by analyzing individual parts in isolation",
                    "Examples include biological ecosystems, economic markets, social networks, and technological infrastructures that all display complex adaptive behaviors",
                    "Understanding complex systems requires interdisciplinary approaches combining mathematics, computer science, physics, biology, and social sciences methodologies"
                ]
            },
            {
                "title": "System Analysis Methodologies", 
                "content": [
                    "Network analysis provides tools for mapping relationships and information flow patterns between different components within complex system architectures",
                    "Agent-based modeling simulates individual entities and their interactions to understand how macro-level patterns emerge from micro-level behaviors and decisions",
                    "Statistical analysis techniques help identify patterns, correlations, and statistical significance in large datasets generated by complex system observations",
                    "Machine learning algorithms can detect hidden patterns and make predictions about system behavior based on historical data and current state information",
                    "Visualization techniques transform complex data into intuitive graphical representations that facilitate human understanding and decision-making processes"
                ]
            },
            {
                "title": "Applications and Future Directions",
                "content": [
                    "Healthcare systems benefit from complex systems analysis through improved patient flow optimization, resource allocation, and epidemic prediction modeling",
                    "Smart city initiatives utilize complex systems principles to optimize traffic flow, energy distribution, waste management, and citizen service delivery",
                    "Financial markets employ complex systems modeling for risk assessment, algorithmic trading strategies, and economic stability analysis and prediction",
                    "Climate science leverages complex systems approaches to model global weather patterns, climate change impacts, and environmental sustainability strategies",
                    "Artificial intelligence development increasingly incorporates complex systems principles to create more robust, adaptive, and intelligent autonomous systems"
                ]
            }
        ],
        "call_to_action": "Questions, Discussion, and Further Research Welcome!"
    }
    
    # Test with template support
    print("🎨 Testing Enhanced PPT Generator with Template Support")
    
    try:
        # Check for available templates
        available_templates = []
        if os.path.exists(templates_folder):
            for file in os.listdir(templates_folder):
                if file.lower().endswith('.pptx') and not file.startswith('~'):
                    available_templates.append(file)
        
        if available_templates:
            print(f"Available templates: {available_templates}")
            template_to_use = available_templates[0]  # Use first available template
        else:
            print("No templates found. Using default PowerPoint template.")
            template_to_use = "default"
        
        generator = EnhancedPPTGenerator(template_to_use)
        
        # Print template information
        template_info = generator.get_template_info()
        print(f"Template Info: {template_info}")
        
        presentation, slide_count = generator.generate_from_content(sample_content)
        filename = generator.save("template_test_presentation")
        
        print(f"✅ Enhanced presentation created successfully!")
        print(f"📊 Total slides: {slide_count}")
        print(f"💾 Saved as: {filename}")
        print(f"🎯 Template used: {template_to_use}")
        print(f"🎯 Features:")
        print(f"   ✓ Template-based design (colors, fonts, backgrounds from template)")
        print(f"   ✓ Comprehensive bullet points (15-35 words each)")
        print(f"   ✓ Enhanced text formatting")
        print(f"   ✓ Improved content distribution and spacing")
        print(f"   ✓ General subject adaptability")
        
    except Exception as e:
        print(f"❌ Error: {e}")
    
    print(f"\n🎯 Key Changes Made:")
    print(f"   ✓ Replaced theme color system with template support")
    print(f"   ✓ Templates loaded from 'templates' folder")
    print(f"   ✓ Template layouts and styling preserved")
    print(f"   ✓ Fallback to default PowerPoint template if template not found")
    print(f"   ✓ Template information tracking and reporting")